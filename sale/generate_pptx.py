"""Generate the Analize.Online Romanian sales pitch PowerPoint.

Visually stunning dark-theme presentation with 14 slides.
Features: decorative background elements, gradient accent bars, slide transitions,
large decorative numbers, visual flow diagrams, icon shapes, glow effects,
card shadows, dramatic title/CTA slides, and progress indicators.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from lxml import etree

# ── Brand colors ──────────────────────────────────────────────
DARK_BG = RGBColor(0x0B, 0x11, 0x1E)
CARD_BG = RGBColor(0x15, 0x1D, 0x2E)
CARD_BG_LIGHT = RGBColor(0x1A, 0x24, 0x38)
CARD_SHADOW = RGBColor(0x08, 0x0D, 0x18)
CYAN = RGBColor(0x06, 0xB6, 0xD4)
CYAN_DIM = RGBColor(0x04, 0x7A, 0x8E)
TEAL = RGBColor(0x14, 0xB8, 0xA6)
TEAL_DIM = RGBColor(0x0D, 0x7A, 0x6F)
WHITE = RGBColor(0xF1, 0xF5, 0xF9)
LIGHT_GRAY = RGBColor(0x94, 0xA3, 0xB8)
MID_GRAY = RGBColor(0x64, 0x74, 0x8B)
DARK_GRAY = RGBColor(0x47, 0x55, 0x69)
GREEN = RGBColor(0x22, 0xC5, 0x5E)
GREEN_DIM = RGBColor(0x16, 0x83, 0x3E)
AMBER = RGBColor(0xFB, 0xBF, 0x24)
AMBER_DIM = RGBColor(0xA8, 0x7F, 0x18)
RED = RGBColor(0xEF, 0x44, 0x44)
VIOLET = RGBColor(0xA7, 0x8B, 0xFA)
VIOLET_DIM = RGBColor(0x6F, 0x5D, 0xA7)

# Glow colors (very dark, subtle)
GLOW_CYAN = RGBColor(0x06, 0x2A, 0x35)
GLOW_TEAL = RGBColor(0x06, 0x2E, 0x2A)
GLOW_VIOLET = RGBColor(0x1A, 0x15, 0x30)
GLOW_GREEN = RGBColor(0x08, 0x28, 0x15)
GLOW_AMBER = RGBColor(0x28, 0x20, 0x08)

# ── Slide dimensions (16:9) ──────────────────────────────────
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ── Layout constants ─────────────────────────────────────────
MARGIN_X = Inches(0.9)
CONTENT_W = Inches(11.5)  # 13.333 - 2*0.9
TAG_TOP = Inches(0.95)
TITLE_TOP = Inches(1.55)
BODY_TOP = Inches(2.75)
FOOTER_TOP = Inches(7.05)
FOOTER_H = Inches(0.35)

# ── Typography sizes ─────────────────────────────────────────
FONT_TITLE = 38
FONT_HEADING = 22
FONT_BODY = 15
FONT_CAPTION = 11
FONT_HERO = 72  # For large decorative numbers


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# HELPER FUNCTIONS
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━


def set_slide_bg(slide, color=DARK_BG):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_slide_transition(slide, transition_type="fade", duration=700):
    """Add a transition to a slide via XML."""
    sld = slide._element
    transition = etree.SubElement(sld, qn('p:transition'))
    transition.set('spd', 'med')
    transition.set('advTm', str(duration))
    if transition_type == "fade":
        etree.SubElement(transition, qn('p:fade'))
    elif transition_type == "push":
        push = etree.SubElement(transition, qn('p:push'))
        push.set('dir', 'l')
    elif transition_type == "wipe":
        wipe = etree.SubElement(transition, qn('p:wipe'))
        wipe.set('dir', 'r')
    elif transition_type == "cover":
        cover = etree.SubElement(transition, qn('p:cover'))
        cover.set('dir', 'r')


def add_text_box(slide, left, top, width, height, text, font_size=FONT_BODY,
                 color=WHITE, bold=False, alignment=PP_ALIGN.LEFT,
                 font_name="Calibri", anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.margin_left = Pt(4)
    tf.margin_right = Pt(4)
    tf.margin_top = Pt(3)
    tf.margin_bottom = Pt(3)
    tf.auto_size = None
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = alignment
    return tf


def add_paragraph(tf, text, font_size=FONT_BODY, color=WHITE, bold=False,
                  alignment=PP_ALIGN.LEFT, space_before=Pt(6), space_after=Pt(2)):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Calibri"
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    return p


def add_rounded_rect(slide, left, top, width, height, fill_color=CARD_BG,
                     border_color=None, border_width=Pt(0)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_card_with_shadow(slide, left, top, width, height, fill_color=CARD_BG,
                         border_color=None, border_width=Pt(0), shadow_offset=Inches(0.06)):
    """Add a card with a shadow effect (darker rect offset behind it)."""
    # Shadow rectangle (slightly offset down-right)
    shadow = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                    left + shadow_offset, top + shadow_offset,
                                    width, height)
    shadow.fill.solid()
    shadow.fill.fore_color.rgb = CARD_SHADOW
    shadow.line.fill.background()
    shadow.shadow.inherit = False

    # Main card on top
    card = add_rounded_rect(slide, left, top, width, height, fill_color,
                            border_color, border_width)
    return card


def add_accent_line(slide, left, top, width, color=CYAN, thickness=Pt(2)):
    """Add a thin horizontal accent line."""
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, thickness)
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    return line


def add_gradient_bar(slide, left, top, width, color1, color2, color3=None,
                     bar_height=Pt(5)):
    """Create a gradient-like bar by stacking thin colored rectangles."""
    h = bar_height // 3 if color3 else bar_height // 2
    r1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, h)
    r1.fill.solid()
    r1.fill.fore_color.rgb = color1
    r1.line.fill.background()

    r2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top + h, width, h)
    r2.fill.solid()
    r2.fill.fore_color.rgb = color2
    r2.line.fill.background()

    if color3:
        r3 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top + h * 2, width, h)
        r3.fill.solid()
        r3.fill.fore_color.rgb = color3
        r3.line.fill.background()


def add_glow(slide, center_x, center_y, radius, color):
    """Add a large semi-transparent circle as a background glow effect."""
    left = center_x - radius
    top = center_y - radius
    size = radius * 2
    glow = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    glow.fill.solid()
    glow.fill.fore_color.rgb = color
    glow.line.fill.background()
    glow.shadow.inherit = False
    return glow


def add_decorative_circle(slide, left, top, size, color, opacity_level=1):
    """Add a decorative semi-transparent circle. opacity_level: 1=very subtle, 2=subtle, 3=visible."""
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    circ.fill.solid()
    # Use very dark tinted colors to simulate transparency on dark bg
    c = str(color)
    r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
    factor = {1: 15, 2: 10, 3: 7}.get(opacity_level, 12)
    blended = RGBColor(
        max(0x0B, min(255, 0x0B + r // factor)),
        max(0x11, min(255, 0x11 + g // factor)),
        max(0x1E, min(255, 0x1E + b // factor))
    )
    circ.fill.fore_color.rgb = blended
    circ.line.fill.background()
    circ.shadow.inherit = False
    return circ


def add_diamond_bullet(slide, left, top, size=Inches(0.12), color=CYAN):
    """Add a small diamond shape as a decorative bullet marker."""
    d = slide.shapes.add_shape(MSO_SHAPE.DIAMOND, left, top, size, size)
    d.fill.solid()
    d.fill.fore_color.rgb = color
    d.line.fill.background()
    d.shadow.inherit = False
    return d


def add_tag(slide, left, top, text, color=CYAN):
    """Small pill-shaped tag label."""
    w, h = Inches(2.6), Inches(0.32)
    c = str(color)
    r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
    bg = RGBColor(r // 8, g // 8, b // 8)
    shape = add_rounded_rect(slide, left, top, w, h, fill_color=bg)
    shape.text_frame.word_wrap = False
    p = shape.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(FONT_CAPTION)
    p.font.color.rgb = color
    p.font.bold = True
    p.font.name = "Calibri"
    p.alignment = PP_ALIGN.CENTER
    shape.text_frame.margin_top = Pt(3)
    shape.text_frame.margin_bottom = Pt(3)


def add_footer(slide, num, total=14):
    """Consistent bottom bar: left text + right slide number."""
    add_accent_line(slide, MARGIN_X, FOOTER_TOP, CONTENT_W,
                    color=RGBColor(0x1E, 0x29, 0x3B), thickness=Pt(1))
    add_text_box(slide, MARGIN_X, FOOTER_TOP + Pt(4), Inches(6), Inches(0.3),
                 "analize.online   |   Confiden\u021bial   |   Martie 2026",
                 font_size=9, color=DARK_GRAY, alignment=PP_ALIGN.LEFT)
    add_text_box(slide, SLIDE_W - Inches(1.8), FOOTER_TOP + Pt(4), Inches(1.2), Inches(0.3),
                 f"{num} / {total}", font_size=9, color=DARK_GRAY, alignment=PP_ALIGN.RIGHT)


def logo_bar(slide):
    add_text_box(slide, MARGIN_X, Inches(0.3), Inches(3), Inches(0.35),
                 "Analize.Online", font_size=13, color=MID_GRAY, bold=True)
    dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, MARGIN_X - Inches(0.25), Inches(0.37),
                                 Inches(0.12), Inches(0.12))
    dot.fill.solid()
    dot.fill.fore_color.rgb = TEAL
    dot.line.fill.background()


def add_stat_badge(slide, left, top, value, label, accent_color=CYAN,
                   badge_w=Inches(2.6), badge_h=Inches(1.8), hero=False):
    """A stat with large number over a card with shadow."""
    card = add_card_with_shadow(slide, left, top, badge_w, badge_h, CARD_BG)

    # Gradient bar at top of card
    c = str(accent_color)
    r, g, b = int(c[0:2], 16), int(c[2:4], 16), int(c[4:6], 16)
    dim_color = RGBColor(r // 2, g // 2, b // 2)
    add_gradient_bar(slide, left + Inches(0.25), top + Inches(0.06),
                     badge_w - Inches(0.5), accent_color, dim_color, bar_height=Pt(4))

    # Large value
    val_size = FONT_HERO if hero else 42
    add_text_box(slide, left + Inches(0.15), top + Inches(0.25),
                 badge_w - Inches(0.3), Inches(0.85),
                 value, font_size=val_size, color=accent_color, bold=True,
                 alignment=PP_ALIGN.CENTER)
    # Label below
    label_top = top + Inches(1.15) if hero else top + Inches(1.1)
    add_text_box(slide, left + Inches(0.15), label_top,
                 badge_w - Inches(0.3), Inches(0.6),
                 label, font_size=FONT_BODY, color=LIGHT_GRAY,
                 alignment=PP_ALIGN.CENTER)
    return card


def add_arrow_connector(slide, x1, y, x2, color=TEAL):
    """Draw a horizontal arrow-like connector between two points."""
    w = x2 - x1
    h = Pt(3)
    line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x1, y, w, h)
    line.fill.solid()
    line.fill.fore_color.rgb = color
    line.line.fill.background()
    # Arrow head (triangle)
    tri_size = Inches(0.15)
    tri = slide.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                                 x2 - Inches(0.02), y - Inches(0.06),
                                 tri_size, tri_size)
    tri.fill.solid()
    tri.fill.fore_color.rgb = color
    tri.line.fill.background()
    # Rotate triangle 90 degrees to point right
    tri.rotation = 90.0


# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 1: Title
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
add_slide_transition(sl, "fade", 800)

# Background decorative elements - large translucent circles
add_decorative_circle(sl, Inches(-2), Inches(-2), Inches(7), CYAN, 1)
add_decorative_circle(sl, Inches(9), Inches(4), Inches(6), TEAL, 1)
add_decorative_circle(sl, Inches(5), Inches(-3), Inches(5), VIOLET, 1)

# Glow effects
add_glow(sl, Inches(6.666), Inches(3.75), Inches(3), GLOW_CYAN)

# Accent lines above and below title
add_gradient_bar(sl, Inches(4.0), Inches(1.9), Inches(5.3), TEAL, CYAN, TEAL_DIM, Pt(5))
add_gradient_bar(sl, Inches(4.0), Inches(4.65), Inches(5.3), TEAL_DIM, CYAN, TEAL, Pt(5))

# Large decorative circle behind title
add_decorative_circle(sl, Inches(3.5), Inches(0.8), Inches(6.3), TEAL, 2)

add_text_box(sl, Inches(1), Inches(2.15), Inches(11.3), Inches(1.0),
             "Analize.Online", font_size=58, color=WHITE, bold=True,
             alignment=PP_ALIGN.CENTER)

add_text_box(sl, Inches(2), Inches(3.3), Inches(9.3), Inches(1.2),
             "Prima platform\u0103 din Rom\u00e2nia care unific\u0103\ntoate analizele medicale cu inteligen\u021b\u0103 artificial\u0103",
             font_size=24, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(sl, Inches(2), Inches(5.0), Inches(9.3), Inches(0.4),
             "Pitch de v\u00e2nzare / licen\u021biere \u2014 Martie 2026",
             font_size=FONT_BODY, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

add_text_box(sl, Inches(2), Inches(5.8), Inches(9.3), Inches(0.4),
             "Confiden\u021bial   \u2022   contact@analize.online   \u2022   analize.online",
             font_size=FONT_CAPTION, color=DARK_GRAY, alignment=PP_ALIGN.CENTER)

add_footer(sl, 1)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 2: Problema
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
add_slide_transition(sl, "push", 700)
logo_bar(sl)

# Background decorative elements
add_decorative_circle(sl, Inches(10), Inches(-1), Inches(5), CYAN, 1)
add_glow(sl, Inches(10.6), Inches(3.5), Inches(2.5), GLOW_CYAN)

add_tag(sl, MARGIN_X, TAG_TOP, "PROBLEMA", CYAN)

add_text_box(sl, MARGIN_X, TITLE_TOP, Inches(7), Inches(1.0),
             "Datele medicale din Rom\u00e2nia\nsunt fragmentate \u0219i inaccesibile",
             font_size=FONT_TITLE, color=WHITE, bold=True)

# Gradient bar under title
add_gradient_bar(sl, MARGIN_X, Inches(2.55), Inches(3.5), CYAN, CYAN_DIM, bar_height=Pt(4))

bullets = [
    "Rezultatele sunt blocate \u00een portaluri separate (Regina Maria, Synevo, MedLife, Sanador)",
    "Nicio vedere unificat\u0103 \u2014 pacien\u021bii uit\u0103 rezultatele anterioare",
    "PDF-uri cu sute de valori \u2014 imposibil de interpretat f\u0103r\u0103 medic",
    "Zero urm\u0103rire tendin\u021be \u2014 colesterolul a crescut sau sc\u0103zut?",
    "Vizite reactive, nu preventive \u2014 pacien\u021bii merg doar c\u00e2nd sunt bolnavi",
]
y_bullet = Inches(2.85)
for b in bullets:
    add_diamond_bullet(sl, MARGIN_X + Inches(0.05), y_bullet + Inches(0.04), Inches(0.1), CYAN)
    add_text_box(sl, MARGIN_X + Inches(0.3), y_bullet, Inches(6.5), Inches(0.35),
                 b, font_size=FONT_BODY, color=LIGHT_GRAY)
    y_bullet += Inches(0.55)

# Stats card (right side) - with shadow
card = add_card_with_shadow(sl, Inches(8.6), Inches(1.8), Inches(4.0), Inches(4.5), CARD_BG)
add_gradient_bar(sl, Inches(9.3), Inches(1.95), Inches(2.6), CYAN, CYAN_DIM, bar_height=Pt(4))

# HUGE decorative number
add_text_box(sl, Inches(8.8), Inches(2.15), Inches(3.6), Inches(1.2), "19M+",
             font_size=FONT_HERO, color=CYAN, bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(sl, Inches(8.8), Inches(3.4), Inches(3.6), Inches(0.4), "rom\u00e2ni",
             font_size=FONT_HEADING, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)

# Small decorative circle behind stat
add_decorative_circle(sl, Inches(9.5), Inches(4.3), Inches(2), CYAN, 2)

stats_lines = "4 furnizori mari de analize\n0 platforme unificate\nMilioane de rezultate \u00een silozuri"
add_text_box(sl, Inches(8.8), Inches(4.2), Inches(3.6), Inches(1.8), stats_lines,
             font_size=FONT_BODY, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

add_footer(sl, 2)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 3: Solutia - Visual pipeline
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
add_slide_transition(sl, "wipe", 700)
logo_bar(sl)

# Background glows
add_glow(sl, Inches(1), Inches(5), Inches(2.5), GLOW_TEAL)
add_decorative_circle(sl, Inches(11), Inches(0), Inches(4), TEAL, 1)

add_tag(sl, MARGIN_X, TAG_TOP, "SOLU\u021aIA", TEAL)

add_text_box(sl, MARGIN_X, TITLE_TOP, Inches(11), Inches(0.7),
             "Un singur loc. Toate analizele. AI care le explic\u0103.",
             font_size=FONT_TITLE, color=WHITE, bold=True)

add_gradient_bar(sl, MARGIN_X, Inches(2.25), Inches(3.5), TEAL, TEAL_DIM, bar_height=Pt(4))

steps = [
    ("1", "Conectezi conturile medicale",
     "Regina Maria, Synevo, MedLife, Sanador \u2014 credentiale criptate AES-256"),
    ("2", "Sincronizare automat\u0103",
     "Crawlere Playwright descarc\u0103 toate PDF-urile, gestionare CAPTCHA inclus\u0103"),
    ("3", "Extrac\u021bie AI",
     "GPT-4o extrage 150+ tipuri de biomarkeri: nume, valoare, unitate, referin\u021b\u0103"),
    ("4", "Speciali\u0219ti AI dinamici",
     "Generalistul analizeaz\u0103, apoi recomand\u0103 speciali\u0219tii relevan\u021bi \u2014 orice specialitate, automat"),
    ("5", "Grafice \u0219i alerte",
     "Evolu\u021bie \u00een timp, valori anormale eviden\u021biate, screening-uri preventive"),
]

y = BODY_TOP
step_spacing = Inches(0.88)
box_left = MARGIN_X + Inches(0.7)
box_w = Inches(10.5)

for i, (num, title, desc) in enumerate(steps):
    # Step box background with shadow
    add_card_with_shadow(sl, box_left, y - Inches(0.05), box_w, Inches(0.75), CARD_BG,
                         shadow_offset=Inches(0.04))

    # Number circle with gradient effect (two overlapping circles)
    circ_bg = sl.shapes.add_shape(MSO_SHAPE.OVAL, MARGIN_X - Inches(0.02),
                                  y - Inches(0.02), Inches(0.49), Inches(0.49))
    circ_bg.fill.solid()
    circ_bg.fill.fore_color.rgb = TEAL_DIM
    circ_bg.line.fill.background()

    circ = sl.shapes.add_shape(MSO_SHAPE.OVAL, MARGIN_X, y, Inches(0.45), Inches(0.45))
    circ.fill.solid()
    circ.fill.fore_color.rgb = TEAL
    circ.line.fill.background()
    p = circ.text_frame.paragraphs[0]
    p.text = num
    p.font.size = Pt(16)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    circ.text_frame.margin_top = Pt(2)

    # Connector line between steps (vertical)
    if i < len(steps) - 1:
        conn_top = y + Inches(0.45)
        conn_h = step_spacing - Inches(0.45)
        conn = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   MARGIN_X + Inches(0.19), conn_top,
                                   Pt(3), conn_h)
        conn.fill.solid()
        conn.fill.fore_color.rgb = TEAL_DIM
        conn.line.fill.background()

    add_text_box(sl, box_left + Inches(0.15), y - Inches(0.02), Inches(4.5), Inches(0.35),
                 title, font_size=18, color=WHITE, bold=True)
    add_text_box(sl, box_left + Inches(0.15), y + Inches(0.32), Inches(9.8), Inches(0.4),
                 desc, font_size=14, color=LIGHT_GRAY)
    y += step_spacing

add_footer(sl, 3)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 4: Ce include platforma
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
add_slide_transition(sl, "fade", 700)
logo_bar(sl)

# Background decorative elements
add_decorative_circle(sl, Inches(-1), Inches(4), Inches(4), VIOLET, 1)
add_glow(sl, Inches(12), Inches(1), Inches(2), GLOW_VIOLET)

add_tag(sl, MARGIN_X, TAG_TOP, "INVENTAR", VIOLET)

add_text_box(sl, MARGIN_X, TITLE_TOP, Inches(11), Inches(0.7),
             "Ce include platforma (cuantificat)",
             font_size=FONT_TITLE, color=WHITE, bold=True)

add_gradient_bar(sl, MARGIN_X, Inches(2.25), Inches(3.5), VIOLET, VIOLET_DIM, bar_height=Pt(4))

items_left = [
    ("28", "Servicii backend", "~11.000 linii Python"),
    ("50+", "Endpoint-uri API", "16 routere"),
    ("30+", "Modele baz\u0103 de date", "PostgreSQL"),
    ("20+", "Pagini frontend", "React + Tailwind"),
    ("4", "Crawlere furnizori", "Playwright + CAPTCHA"),
    ("6", "Agen\u021bi AI", "GPT-4o speciali\u0219ti"),
]
items_right = [
    ("150+", "Tipuri biomarkeri", "Normaliza\u021bi"),
    ("2", "Limbi", "RO + EN (i18next)"),
    ("AES-256", "Criptare per user", "Vault + recovery key"),
    ("GDPR", "Conformitate", "Export, \u0219tergere, consent"),
    ("Netopia", "Pl\u0103\u021bi integrate", "29/49 RON/lun\u0103"),
    ("Prometheus", "Monitorizare", "Grafana + Sentry"),
]

y_start = BODY_TOP
row_spacing = Inches(0.6)
for col_idx, items in enumerate([items_left, items_right]):
    x_base = MARGIN_X if col_idx == 0 else Inches(7.0)
    y = y_start
    for val, label, detail in items:
        # Row background with subtle shadow
        add_card_with_shadow(sl, x_base - Inches(0.1), y - Inches(0.04),
                             Inches(5.7), Inches(0.48), CARD_BG, shadow_offset=Inches(0.03))
        # Diamond bullet
        add_diamond_bullet(sl, x_base + Inches(0.05), y + Inches(0.1), Inches(0.08), VIOLET)
        add_text_box(sl, x_base + Inches(0.25), y, Inches(1.15), Inches(0.4),
                     val, font_size=17, color=VIOLET, bold=True)
        add_text_box(sl, x_base + Inches(1.5), y, Inches(2.3), Inches(0.4),
                     label, font_size=FONT_BODY, color=WHITE, bold=True)
        add_text_box(sl, x_base + Inches(3.9), y, Inches(1.8), Inches(0.4),
                     detail, font_size=13, color=MID_GRAY)
        y += row_spacing

add_footer(sl, 4)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 5: Crawlere - Moat tehnic
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
add_slide_transition(sl, "cover", 700)
logo_bar(sl)

# Background elements
add_decorative_circle(sl, Inches(8), Inches(-2), Inches(5), CYAN, 1)
add_glow(sl, Inches(2), Inches(6), Inches(2), GLOW_CYAN)

add_tag(sl, MARGIN_X, TAG_TOP, "MOAT TEHNIC", CYAN)

add_text_box(sl, MARGIN_X, TITLE_TOP, Inches(11), Inches(0.7),
             "4 crawlere. Cel mai greu de replicat.",
             font_size=FONT_TITLE, color=WHITE, bold=True)

add_gradient_bar(sl, MARGIN_X, Inches(2.25), Inches(3.5), CYAN, CYAN_DIM, bar_height=Pt(4))

providers = [
    ("Regina Maria", "CAPTCHA handling\nVisible browser fallback"),
    ("Synevo", "Full headless mode\nMulti-page navigation"),
    ("MedLife", "CAPTCHA detection\nAuto-retry logic"),
    ("Sanador", "Headless automation\nSession management"),
]
x = MARGIN_X
card_w = Inches(2.7)
card_gap = Inches(0.25)
for name, details in providers:
    card = add_card_with_shadow(sl, x, BODY_TOP, card_w, Inches(1.7), CARD_BG)
    add_gradient_bar(sl, x + Inches(0.35), BODY_TOP + Inches(0.06),
                     card_w - Inches(0.7), CYAN, CYAN_DIM, bar_height=Pt(3))
    # Icon circle
    icon_size = Inches(0.3)
    icon = sl.shapes.add_shape(MSO_SHAPE.OVAL,
                               x + (card_w - icon_size) // 2, BODY_TOP + Inches(0.2),
                               icon_size, icon_size)
    icon.fill.solid()
    icon.fill.fore_color.rgb = CYAN_DIM
    icon.line.fill.background()
    add_text_box(sl, x + Inches(0.2), BODY_TOP + Inches(0.55), card_w - Inches(0.4), Inches(0.4),
                 name, font_size=18, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, x + Inches(0.2), BODY_TOP + Inches(0.95), card_w - Inches(0.4), Inches(0.65),
                 details, font_size=13, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
    x += card_w + card_gap

# Why it's a moat
moat_top = Inches(4.7)
moat_h = Inches(2.1)
add_card_with_shadow(sl, MARGIN_X, moat_top, CONTENT_W, moat_h, CARD_BG)
tf = add_text_box(sl, MARGIN_X + Inches(0.25), moat_top + Inches(0.15), Inches(11), Inches(0.35),
                  "De ce este o barier\u0103 de intrare:", font_size=17, color=WHITE, bold=True)
moat_bullets = [
    "F\u0103r\u0103 API-uri publice \u2014 automatizare browser obligatorie (Playwright)",
    "Fiecare furnizor = flux unic de login, CAPTCHA, sesiune, navigare",
    "Gestionare concuren\u021b\u0103: max 2 sincroniz\u0103ri simultane, retry, cleanup",
    "Un concurent ar avea nevoie de 6-12 luni doar pentru crawlere",
]
bullet_y = moat_top + Inches(0.55)
for b in moat_bullets:
    add_diamond_bullet(sl, MARGIN_X + Inches(0.35), bullet_y + Inches(0.06), Inches(0.1), CYAN)
    add_paragraph(tf, f"     {b}", font_size=FONT_BODY, color=LIGHT_GRAY, space_before=Pt(8))
    bullet_y += Inches(0.32)

add_footer(sl, 5)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 6: AI Pipeline
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
add_slide_transition(sl, "fade", 700)
logo_bar(sl)

# Background
add_decorative_circle(sl, Inches(9), Inches(3), Inches(5), VIOLET, 1)
add_glow(sl, Inches(3), Inches(5), Inches(2), GLOW_VIOLET)

add_tag(sl, MARGIN_X, TAG_TOP, "INTELIGEN\u021a\u0102 ARTIFICIAL\u0102", VIOLET)

add_text_box(sl, MARGIN_X, TITLE_TOP, Inches(11), Inches(0.7),
             "Speciali\u0219ti AI dinamici. Powered by GPT-4o.",
             font_size=FONT_TITLE, color=WHITE, bold=True)

add_gradient_bar(sl, MARGIN_X, Inches(2.25), Inches(3.5), VIOLET, VIOLET_DIM, bar_height=Pt(4))

# Left column: flow with connected shapes
flow_items = [
    ("1. Generalist", "Evalueaz\u0103 toate rezultatele, identific\u0103 probleme"),
    ("2. Recomand\u0103", "Pe baza rezultatelor, nu o list\u0103 fix\u0103 \u2014 orice specialitate"),
    ("3. Analizeaz\u0103", "Fiecare specialist prime\u0219te biomarkerii relevan\u021bi"),
    ("Exemple:", "Cardiolog, Endocrinolog, Hematolog, Nefrolog..."),
    ("Dinamic:", "Probleme dermatologice \u2192 Dermatolog AI"),
]
y = BODY_TOP
row_h = Inches(0.52)
row_gap = Inches(0.1)
for i, (label, desc) in enumerate(flow_items):
    add_card_with_shadow(sl, MARGIN_X, y, Inches(5.6), row_h, CARD_BG,
                         shadow_offset=Inches(0.03))
    # Circle bullet for flow
    circ_size = Inches(0.18)
    circ = sl.shapes.add_shape(MSO_SHAPE.OVAL, MARGIN_X + Inches(0.08),
                               y + Inches(0.17), circ_size, circ_size)
    circ.fill.solid()
    circ.fill.fore_color.rgb = VIOLET if i < 3 else TEAL
    circ.line.fill.background()

    add_text_box(sl, MARGIN_X + Inches(0.35), y + Inches(0.07), Inches(1.65), Inches(0.38),
                 label, font_size=14, color=TEAL, bold=True)
    add_text_box(sl, MARGIN_X + Inches(2.0), y + Inches(0.07), Inches(3.4), Inches(0.38),
                 desc, font_size=13, color=LIGHT_GRAY)

    # Vertical connector between flow items
    if i < len(flow_items) - 1:
        conn_top = y + row_h
        conn = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   MARGIN_X + Inches(0.15), conn_top,
                                   Pt(2), row_gap)
        conn.fill.solid()
        conn.fill.fore_color.rgb = VIOLET_DIM
        conn.line.fill.background()

    y += row_h + row_gap

# Right column: features card with shadow
features_top = BODY_TOP
features_h = Inches(3.8)
card = add_card_with_shadow(sl, Inches(7.0), features_top, Inches(5.6), features_h, CARD_BG)
add_gradient_bar(sl, Inches(7.35), features_top + Inches(0.06),
                 Inches(4.9), VIOLET, VIOLET_DIM, bar_height=Pt(3))
tf = add_text_box(sl, Inches(7.3), features_top + Inches(0.2), Inches(5.0), Inches(3.4),
                  "Func\u021bionalit\u0103\u021bi AI pipeline:", font_size=17, color=WHITE, bold=True)
features = [
    "Extrage 150+ tipuri biomarkeri din PDF-uri rom\u00e2ne\u0219ti",
    "Normalizeaz\u0103 denumiri \u00eentre furnizori",
    "Detecteaz\u0103 tendin\u021be pe date istorice",
    "Output bilingv (Rom\u00e2n\u0103 + Englez\u0103)",
    "Recomand\u0103ri lifestyle bazate pe rezultate",
    "Screening-uri preventive recomandate",
]
feat_y = features_top + Inches(0.65)
for f in features:
    add_diamond_bullet(sl, Inches(7.5), feat_y + Inches(0.05), Inches(0.1), VIOLET)
    add_paragraph(tf, f"     {f}", font_size=14, color=LIGHT_GRAY, space_before=Pt(8))
    feat_y += Inches(0.4)

add_footer(sl, 6)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 7: Securitate & Conformitate
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
add_slide_transition(sl, "push", 700)
logo_bar(sl)

# Background
add_decorative_circle(sl, Inches(-2), Inches(2), Inches(5), GREEN, 1)
add_glow(sl, Inches(11), Inches(5), Inches(2), GLOW_GREEN)

add_tag(sl, MARGIN_X, TAG_TOP, "SECURITATE", GREEN)

add_text_box(sl, MARGIN_X, TITLE_TOP, Inches(11.5), Inches(0.7),
             "Criptare per utilizator. Nici admin-ul nu vede datele.",
             font_size=FONT_TITLE, color=WHITE, bold=True)

add_gradient_bar(sl, MARGIN_X, Inches(2.25), Inches(3.5), GREEN, GREEN_DIM, bar_height=Pt(4))

sec_cards = [
    ("Vault per utilizator", "AES-256-GCM",
     "Fiecare user are propria cheie\nderivat\u0103 din parol\u0103\n(PBKDF2, 600K itera\u021bii)"),
    ("Recovery key", "Generat la \u00eenregistrare",
     "Singura metod\u0103 de recuperare\ndac\u0103 parola e pierdut\u0103.\nAdmin-ul NU poate ajuta."),
    ("GDPR complet", "Export + \u0218tergere",
     "Export toate datele (JSON)\n\u0218tergere cont permanent\u0103\nConsim\u021b\u0103m\u00e2nt cookies"),
    ("Autentificare", "JWT + Google OAuth",
     "Rate limiting pe auth\nAudit logging\nHTTPS + security headers"),
]
card_w = Inches(2.7)
card_gap = Inches(0.2)
card_h = Inches(3.5)
x = MARGIN_X
for title, subtitle, desc in sec_cards:
    card = add_card_with_shadow(sl, x, BODY_TOP, card_w, card_h, CARD_BG)
    add_gradient_bar(sl, x + Inches(0.3), BODY_TOP + Inches(0.06),
                     card_w - Inches(0.6), GREEN, GREEN_DIM, bar_height=Pt(3))
    # Shield icon (circle with colored fill)
    shield_size = Inches(0.3)
    shield = sl.shapes.add_shape(MSO_SHAPE.OVAL,
                                 x + (card_w - shield_size) // 2, BODY_TOP + Inches(0.2),
                                 shield_size, shield_size)
    shield.fill.solid()
    shield.fill.fore_color.rgb = GREEN_DIM
    shield.line.fill.background()

    add_text_box(sl, x + Inches(0.2), BODY_TOP + Inches(0.6), card_w - Inches(0.4), Inches(0.4),
                 title, font_size=16, color=GREEN, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, x + Inches(0.2), BODY_TOP + Inches(1.0), card_w - Inches(0.4), Inches(0.4),
                 subtitle, font_size=13, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, x + Inches(0.2), BODY_TOP + Inches(1.6), card_w - Inches(0.4), Inches(1.7),
                 desc, font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
    x += card_w + card_gap

add_footer(sl, 7)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 8: Avantaj competitiv - 3 perspective
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
add_slide_transition(sl, "wipe", 700)
logo_bar(sl)

# Background
add_decorative_circle(sl, Inches(10), Inches(-1), Inches(4), AMBER, 1)
add_decorative_circle(sl, Inches(-1), Inches(5), Inches(4), CYAN, 1)
add_glow(sl, Inches(6.5), Inches(4), Inches(2.5), GLOW_AMBER)

add_tag(sl, MARGIN_X, TAG_TOP, "AVANTAJ COMPETITIV", AMBER)

add_text_box(sl, MARGIN_X, TITLE_TOP, Inches(11), Inches(0.7),
             "Trei perspective de expert",
             font_size=FONT_TITLE, color=WHITE, bold=True)

add_gradient_bar(sl, MARGIN_X, Inches(2.25), Inches(3.5), AMBER, AMBER_DIM, bar_height=Pt(4))

perspectives = [
    ("Expert SaaS", CYAN,
     "First-mover pe o pia\u021b\u0103 f\u0103r\u0103\nconcuren\u021bi. Crawlerele sunt\nbariera de intrare \u2014 6-12 luni\nde replicat. Time-to-market\nsavings: 12+ luni \u0219i \u20ac50-150K."),
    ("Expert S\u0103n\u0103tate Digital\u0103", TEAL,
     "Directiva EU EHDS impune\ndigitalizare date medicale.\nCorporate wellness \u00een explozie.\nAsigur\u0103torii caut\u0103 date\nlongitudinale pentru pricing."),
    ("Expert Evaluare Startup", AMBER,
     "Build vs Buy: 2-3 seniori \u00d7\n12 luni = \u20ac72K-180K + risc.\nCu achizi\u021bia: produs func\u021bional\nimediat, f\u0103r\u0103 risc de e\u0219ec.\nValoarea e \u00een execu\u021bie."),
]
card_w = Inches(3.7)
card_gap = Inches(0.25)
card_h = Inches(3.8)
x = MARGIN_X
for title, color, text in perspectives:
    card = add_card_with_shadow(sl, x, BODY_TOP, card_w, card_h, CARD_BG)
    add_gradient_bar(sl, x + Inches(0.35), BODY_TOP + Inches(0.06),
                     card_w - Inches(0.7), color,
                     RGBColor(int(str(color)[0:2], 16) // 2,
                              int(str(color)[2:4], 16) // 2,
                              int(str(color)[4:6], 16) // 2),
                     bar_height=Pt(3))
    # Decorative circle behind title
    add_decorative_circle(sl, x + Inches(0.3), BODY_TOP + Inches(0.15),
                          Inches(0.5), color, 3)
    add_text_box(sl, x + Inches(0.3), BODY_TOP + Inches(0.3), card_w - Inches(0.6), Inches(0.4),
                 title, font_size=19, color=color, bold=True)
    add_text_box(sl, x + Inches(0.3), BODY_TOP + Inches(0.85), card_w - Inches(0.6), Inches(2.7),
                 text, font_size=FONT_BODY, color=LIGHT_GRAY)
    x += card_w + card_gap

add_footer(sl, 8)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 9: Model de pret incremental
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
add_slide_transition(sl, "fade", 700)
logo_bar(sl)

# Background
add_decorative_circle(sl, Inches(5), Inches(-2), Inches(5), AMBER, 1)
add_glow(sl, Inches(1), Inches(6), Inches(2), GLOW_AMBER)

add_tag(sl, MARGIN_X, TAG_TOP, "EVALUARE", AMBER)

add_text_box(sl, MARGIN_X, TITLE_TOP, Inches(11), Inches(0.6),
             "Model de pre\u021b incremental",
             font_size=FONT_TITLE, color=WHITE, bold=True)

add_gradient_bar(sl, MARGIN_X, Inches(2.15), Inches(3.5), AMBER, AMBER_DIM, bar_height=Pt(4))

add_text_box(sl, MARGIN_X, Inches(2.35), Inches(11), Inches(0.4),
             "Pre\u021bul cre\u0219te cu fiecare etap\u0103 de maturitate. Platforma se afl\u0103 \u00een Etapa 2.",
             font_size=FONT_BODY, color=LIGHT_GRAY)

# Progress arrow/timeline connecting the 4 stages
timeline_y = Inches(2.75)
timeline_colors = [MID_GRAY, TEAL, CYAN, AMBER]
stage_positions = []
card_w = Inches(2.75)
card_gap = Inches(0.15)

# Draw timeline bar
total_timeline_w = 4 * card_w + 3 * card_gap
timeline_bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                   MARGIN_X + Inches(0.2), timeline_y,
                                   total_timeline_w - Inches(0.4), Pt(3))
timeline_bar.fill.solid()
timeline_bar.fill.fore_color.rgb = DARK_GRAY
timeline_bar.line.fill.background()

# Timeline dots
x_timeline = MARGIN_X
for i, tc in enumerate(timeline_colors):
    dot_x = x_timeline + card_w // 2 - Inches(0.1)
    dot = sl.shapes.add_shape(MSO_SHAPE.OVAL, dot_x, timeline_y - Inches(0.07),
                              Inches(0.2), Inches(0.2))
    dot.fill.solid()
    dot.fill.fore_color.rgb = tc
    dot.line.fill.background()
    stage_positions.append(x_timeline)
    x_timeline += card_w + card_gap

# Arrow at end of timeline
arrow_x = MARGIN_X + total_timeline_w - Inches(0.3)
arrow = sl.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE,
                            arrow_x, timeline_y - Inches(0.08),
                            Inches(0.2), Inches(0.2))
arrow.fill.solid()
arrow.fill.fore_color.rgb = AMBER
arrow.line.fill.background()
arrow.rotation = 90.0

stages = [
    ("Etapa 1", "Func\u021bional,\nnelansat", "15.000 \u2013\n30.000 \u20ac",
     "Cod surs\u0103, crawlere,\nAI pipeline. 50% risk\ndiscount.", MID_GRAY),
    ("Etapa 2", "Lansat, primii\nutilizatori", "35.000 \u2013\n60.000 \u20ac",
     "Live, 8 users, infra\ncomplet\u0103. 3-6 luni\nsuport tranzi\u021bie.", TEAL),
    ("Etapa 3", "Trac\u021biune\n100+ pl\u0103titori", "60.000 \u2013\n120.000 \u20ac",
     "Venituri recurente,\n2-4x ARR +\nvaloare active.", CYAN),
    ("Etapa 4", "Scale\n1000+ pl\u0103titori", "150.000 \u2013\n300.000 \u20ac",
     "Achizi\u021bie strategic\u0103,\n5-8x ARR, brand\nrecunoscut.", AMBER),
]
card_h = Inches(3.7)
cards_top = Inches(3.1)
x = MARGIN_X
for idx, (name, subtitle, price, desc, color) in enumerate(stages):
    is_current = name == "Etapa 2"
    border_c = TEAL if is_current else None
    border_w = Pt(2.5) if is_current else Pt(0)
    card = add_card_with_shadow(sl, x, cards_top, card_w, card_h, CARD_BG,
                                border_color=border_c, border_width=border_w)

    tag_text = name + ("  \u2190 ACUM" if is_current else "")
    add_text_box(sl, x + Inches(0.15), cards_top + Inches(0.2), card_w - Inches(0.3), Inches(0.3),
                 tag_text, font_size=13, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, x + Inches(0.15), cards_top + Inches(0.55), card_w - Inches(0.3), Inches(0.55),
                 subtitle, font_size=12, color=MID_GRAY, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, x + Inches(0.15), cards_top + Inches(1.2), card_w - Inches(0.3), Inches(0.9),
                 price, font_size=FONT_HEADING, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, x + Inches(0.15), cards_top + Inches(2.2), card_w - Inches(0.3), Inches(1.3),
                 desc, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    x += card_w + card_gap

add_footer(sl, 9)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 10: Build vs Buy
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
add_slide_transition(sl, "push", 700)
logo_bar(sl)

# Background
add_decorative_circle(sl, Inches(0), Inches(-1), Inches(4), CYAN, 1)
add_decorative_circle(sl, Inches(10), Inches(4), Inches(4), VIOLET, 1)
add_glow(sl, Inches(6.5), Inches(4), Inches(2.5), GLOW_TEAL)

add_tag(sl, MARGIN_X, TAG_TOP, "BUILD VS BUY", CYAN)

add_text_box(sl, MARGIN_X, TITLE_TOP, Inches(11), Inches(0.7),
             "De ce s\u0103 cumperi, nu s\u0103 construie\u0219ti?",
             font_size=FONT_TITLE, color=WHITE, bold=True)

add_gradient_bar(sl, MARGIN_X, Inches(2.25), Inches(3.5), CYAN, CYAN_DIM, bar_height=Pt(4))

col_w = Inches(3.6)
col_gap = Inches(0.25)
col_h = Inches(3.9)
col_top = BODY_TOP
cols = [
    (MARGIN_X, "DAC\u0102 CONSTRUIE\u0218TI", "\u20ac72.000 \u2013 180.000", "+ 12 luni minim",
     "2-3 programatori seniori\nReverse-engineering crawlere\nPipeline AI de la zero\nSecuritate & GDPR\nNicio garan\u021bie de succes",
     RED, False),
    (MARGIN_X + col_w + col_gap, "DAC\u0102 CUMPERI", "\u20ac35.000 \u2013 60.000", "Imediat",
     "Produs func\u021bional azi\n4 crawlere deja construite\nPipeline AI opera\u021bional\nGDPR conformant\nUtilizatori reali",
     TEAL, True),
    (MARGIN_X + 2 * (col_w + col_gap), "ALTERNATIV: LICEN\u021a\u0102", "\u20ac500 \u2013 2.000", "Per lun\u0103, per client",
     "White-label sub brand-ul t\u0103u\nDomeniu propriu\nMentenan\u021b\u0103 inclus\u0103\nF\u0103r\u0103 cost ini\u021bial mare\nOp\u021biune revenue share",
     VIOLET, False),
]
for x, header, price, time_label, items, color, bordered in cols:
    border_c = color if bordered else None
    border_w = Pt(2) if bordered else Pt(0)
    card = add_card_with_shadow(sl, x, col_top, col_w, col_h, CARD_BG,
                                border_color=border_c, border_width=border_w)
    add_gradient_bar(sl, x + Inches(0.35), col_top + Inches(0.06),
                     col_w - Inches(0.7), color,
                     RGBColor(int(str(color)[0:2], 16) // 2,
                              int(str(color)[2:4], 16) // 2,
                              int(str(color)[4:6], 16) // 2),
                     bar_height=Pt(3))
    add_text_box(sl, x + Inches(0.2), col_top + Inches(0.25), col_w - Inches(0.4), Inches(0.3),
                 header, font_size=13, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, x + Inches(0.2), col_top + Inches(0.65), col_w - Inches(0.4), Inches(0.6),
                 price, font_size=26, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, x + Inches(0.2), col_top + Inches(1.3), col_w - Inches(0.4), Inches(0.35),
                 time_label, font_size=FONT_BODY, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, x + Inches(0.2), col_top + Inches(1.85), col_w - Inches(0.4), Inches(1.8),
                 items, font_size=13, color=MID_GRAY, alignment=PP_ALIGN.CENTER)

add_footer(sl, 10)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 11: Modele de tranzactie
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
add_slide_transition(sl, "cover", 700)
logo_bar(sl)

# Background
add_decorative_circle(sl, Inches(11), Inches(-1), Inches(5), TEAL, 1)
add_glow(sl, Inches(2), Inches(6), Inches(2), GLOW_TEAL)

add_tag(sl, MARGIN_X, TAG_TOP, "OP\u021aIUNI", TEAL)

add_text_box(sl, MARGIN_X, TITLE_TOP, Inches(11), Inches(0.7),
             "Trei modele de tranzac\u021bie",
             font_size=FONT_TITLE, color=WHITE, bold=True)

add_gradient_bar(sl, MARGIN_X, Inches(2.25), Inches(3.5), TEAL, TEAL_DIM, bar_height=Pt(4))

models = [
    ("A", "V\u00e2nzare de active", "35.000 \u2013 60.000 EUR",
     "Pre\u021b unic. Cod surs\u0103 complet,\ndomeniu, infrastructur\u0103, baz\u0103 de\ndate, documenta\u021bie.\n3-6 luni suport tranzi\u021bie.\n\nIdeal: companii cu echip\u0103\ntehnic\u0103 proprie.",
     CYAN),
    ("B", "Licen\u021b\u0103 White-Label", "500 \u2013 2.000 EUR / lun\u0103",
     "Platform\u0103 sub brand-ul t\u0103u,\ndomeniu propriu. Mentenan\u021b\u0103\ncrawlere, actualiz\u0103ri, hosting,\nsuport tehnic incluse.\n\nIdeal: clinici, asigur\u0103tori,\nwellness.",
     VIOLET),
    ("C", "Parteneriat Revenue Share", "10K\u201320K EUR + 20-40%",
     "Investi\u021bie ini\u021bial\u0103 redus\u0103.\nAcces complet la platform\u0103.\nDezvoltare continu\u0103.\n\nIdeal: investitori care vor\nupside pe termen lung.",
     AMBER),
]
card_w = Inches(3.6)
card_gap = Inches(0.25)
card_h = Inches(4.0)
x = MARGIN_X
for letter, title, price, desc, color in models:
    card = add_card_with_shadow(sl, x, BODY_TOP, card_w, card_h, CARD_BG)
    add_gradient_bar(sl, x + Inches(0.35), BODY_TOP + Inches(0.06),
                     card_w - Inches(0.7), color,
                     RGBColor(int(str(color)[0:2], 16) // 2,
                              int(str(color)[2:4], 16) // 2,
                              int(str(color)[4:6], 16) // 2),
                     bar_height=Pt(3))
    # Letter circle with glow behind it
    circ_size = Inches(0.5)
    circ_x = x + (card_w - circ_size) // 2
    # Glow behind circle
    glow_size = Inches(0.7)
    glow_circ = sl.shapes.add_shape(MSO_SHAPE.OVAL,
                                    circ_x - Inches(0.1), BODY_TOP + Inches(0.15),
                                    glow_size, glow_size)
    glow_circ.fill.solid()
    c_str = str(color)
    glow_circ.fill.fore_color.rgb = RGBColor(
        int(c_str[0:2], 16) // 6, int(c_str[2:4], 16) // 6, int(c_str[4:6], 16) // 6)
    glow_circ.line.fill.background()

    circ = sl.shapes.add_shape(MSO_SHAPE.OVAL, circ_x, BODY_TOP + Inches(0.25),
                               circ_size, circ_size)
    circ.fill.solid()
    circ.fill.fore_color.rgb = color
    circ.line.fill.background()
    p = circ.text_frame.paragraphs[0]
    p.text = letter
    p.font.size = Pt(18)
    p.font.color.rgb = WHITE
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    circ.text_frame.margin_top = Pt(2)

    add_text_box(sl, x + Inches(0.2), BODY_TOP + Inches(0.9), card_w - Inches(0.4), Inches(0.35),
                 title, font_size=17, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, x + Inches(0.2), BODY_TOP + Inches(1.3), card_w - Inches(0.4), Inches(0.35),
                 price, font_size=14, color=color, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_box(sl, x + Inches(0.2), BODY_TOP + Inches(1.8), card_w - Inches(0.4), Inches(2.0),
                 desc, font_size=12, color=LIGHT_GRAY, alignment=PP_ALIGN.CENTER)
    x += card_w + card_gap

add_footer(sl, 11)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 12: Tractiune actuala
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
add_slide_transition(sl, "wipe", 700)
logo_bar(sl)

# Background
add_decorative_circle(sl, Inches(8), Inches(0), Inches(5), GREEN, 1)
add_decorative_circle(sl, Inches(-2), Inches(3), Inches(4), CYAN, 1)
add_glow(sl, Inches(6.5), Inches(3), Inches(2.5), GLOW_GREEN)

add_tag(sl, MARGIN_X, TAG_TOP, "TRAC\u021aIUNE", GREEN)

add_text_box(sl, MARGIN_X, TITLE_TOP, Inches(11), Inches(0.6),
             "Early, dar real",
             font_size=FONT_TITLE, color=WHITE, bold=True)

add_gradient_bar(sl, MARGIN_X, Inches(2.15), Inches(3.5), GREEN, GREEN_DIM, bar_height=Pt(4))

# Stats with HUGE numbers
stats = [
    ("8", "Utilizatori\n\u00eenregistra\u021bi"),
    ("59", "Documente\nprocesate"),
    ("731", "Biomarkeri\nextra\u0219i"),
    ("100%", "Uptime\n(3 luni)"),
]
stat_w = Inches(2.7)
stat_gap = Inches(0.2)
x = MARGIN_X
for val, label in stats:
    add_stat_badge(sl, x, BODY_TOP, val, label, accent_color=CYAN,
                   badge_w=stat_w, badge_h=Inches(1.8), hero=True)
    x += stat_w + stat_gap

# What's included
inc_top = Inches(4.8)
inc_h = Inches(2.0)
card = add_card_with_shadow(sl, MARGIN_X, inc_top, CONTENT_W, inc_h, CARD_BG)
add_text_box(sl, MARGIN_X + Inches(0.3), inc_top + Inches(0.15), Inches(5), Inches(0.35),
             "Ce prime\u0219te cump\u0103r\u0103torul:", font_size=17, color=WHITE, bold=True)

left_items = [
    "\u2713  Cod surs\u0103 complet (frontend + backend)",
    "\u2713  4 crawlere func\u021bionale",
    "\u2713  Agen\u021bi AI speciali\u0219ti dinamici",
    "\u2713  Domeniu: analize.online",
    "\u2713  Certificat SSL (valid Apr 2026)",
    "\u2713  Infrastructur\u0103 VPS produc\u021bie",
]
right_items = [
    "\u2713  Baz\u0103 de date PostgreSQL",
    "\u2713  Monitorizare Prometheus + Grafana",
    "\u2713  Email AWS SES configurat",
    "\u2713  Pl\u0103\u021bi Netopia integrate",
    "\u2713  Documenta\u021bie complet\u0103",
    "\u2713  3-6 luni suport tranzi\u021bie",
]

tf = add_text_box(sl, MARGIN_X + Inches(0.3), inc_top + Inches(0.5), Inches(5.3), Inches(1.4),
                  "", font_size=13, color=LIGHT_GRAY)
tf.paragraphs[0].text = ""
for item in left_items:
    add_paragraph(tf, item, font_size=12, color=LIGHT_GRAY, space_before=Pt(2), space_after=Pt(1))

tf2 = add_text_box(sl, Inches(6.6), inc_top + Inches(0.5), Inches(5.3), Inches(1.4),
                   "", font_size=13, color=LIGHT_GRAY)
tf2.paragraphs[0].text = ""
for item in right_items:
    add_paragraph(tf2, item, font_size=12, color=LIGHT_GRAY, space_before=Pt(2), space_after=Pt(1))

add_footer(sl, 12)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 13: Riscuri (transparenta)
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
add_slide_transition(sl, "fade", 700)
logo_bar(sl)

# Background
add_decorative_circle(sl, Inches(10), Inches(3), Inches(4), AMBER, 1)
add_glow(sl, Inches(1), Inches(1), Inches(2), GLOW_AMBER)

add_tag(sl, MARGIN_X, TAG_TOP, "TRANSPAREN\u021a\u0102", AMBER)

add_text_box(sl, MARGIN_X, TITLE_TOP, Inches(11), Inches(0.6),
             "Riscuri \u0219i cum le adres\u0103m",
             font_size=FONT_TITLE, color=WHITE, bold=True)

add_gradient_bar(sl, MARGIN_X, Inches(2.15), Inches(3.5), AMBER, AMBER_DIM, bar_height=Pt(4))

risks = [
    ("Crawlerele se pot strica", "Medie",
     "Cod modular, documentat. Update tipic: 1-3 zile per crawler"),
    ("Bus factor = 1", "Medie",
     "Cod curat, documenta\u021bie complet\u0103, tranzi\u021bie asistat\u0103 3-6 luni"),
    ("Dependen\u021b\u0103 OpenAI", "Sc\u0103zut\u0103",
     "Pipeline modular \u2014 \u00eenlocuibil cu Claude, Gemini, sau model local"),
    ("GDPR / date medicale", "Sc\u0103zut\u0103",
     "Criptare per user, export date, \u0219tergere cont func\u021bionale"),
    ("Scalare", "Sc\u0103zut\u0103",
     "PostgreSQL, async architecture, concurrency limits"),
]

row_h = Inches(0.7)
row_gap = Inches(0.12)
y = BODY_TOP
for risk, severity, mitigation in risks:
    card = add_card_with_shadow(sl, MARGIN_X, y, CONTENT_W, row_h, CARD_BG,
                                shadow_offset=Inches(0.04))
    sev_color = AMBER if severity == "Medie" else GREEN
    # Severity diamond
    add_diamond_bullet(sl, MARGIN_X + Inches(0.2), y + Inches(0.12), Inches(0.14), sev_color)

    add_text_box(sl, MARGIN_X + Inches(0.5), y + Inches(0.05), Inches(3.2), Inches(0.3),
                 risk, font_size=FONT_BODY, color=WHITE, bold=True)
    add_text_box(sl, MARGIN_X + Inches(0.5), y + Inches(0.38), Inches(1.5), Inches(0.25),
                 severity, font_size=FONT_CAPTION, color=sev_color, bold=True)
    add_text_box(sl, Inches(4.5), y + Inches(0.15), Inches(7.5), Inches(0.4),
                 mitigation, font_size=14, color=LIGHT_GRAY)
    y += row_h + row_gap

add_footer(sl, 13)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SLIDE 14: Contact / CTA
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_slide_bg(sl)
add_slide_transition(sl, "fade", 1000)

# Dramatic background decorative elements
add_decorative_circle(sl, Inches(-3), Inches(-2), Inches(8), TEAL, 1)
add_decorative_circle(sl, Inches(9), Inches(3), Inches(7), CYAN, 1)
add_decorative_circle(sl, Inches(3), Inches(-1), Inches(5), VIOLET, 1)

# Large glow behind CTA
add_glow(sl, Inches(6.666), Inches(3.75), Inches(4), GLOW_TEAL)
add_glow(sl, Inches(6.666), Inches(3.75), Inches(2.5), GLOW_CYAN)

# Accent lines
add_gradient_bar(sl, Inches(3.5), Inches(1.3), Inches(6.3), TEAL, CYAN, TEAL_DIM, Pt(5))

add_text_box(sl, Inches(1), Inches(1.6), Inches(11.3), Inches(1.2),
             "Programeaz\u0103 un demo live\nde 15 minute.",
             font_size=44, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER)

# CTA card - dramatic with gradient-like background
cta_w = Inches(7.5)
cta_h = Inches(3.0)
cta_left = (SLIDE_W - cta_w) // 2

# Shadow behind CTA
shadow = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                             cta_left + Inches(0.1), Inches(3.25),
                             cta_w, cta_h)
shadow.fill.solid()
shadow.fill.fore_color.rgb = RGBColor(0x08, 0x3D, 0x35)
shadow.line.fill.background()

# Main CTA card
cta = add_rounded_rect(sl, cta_left, Inches(3.15), cta_w, cta_h, TEAL)

# Gradient-like overlay strips on CTA
overlay1 = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               cta_left, Inches(3.15), cta_w, Inches(0.15))
overlay1.fill.solid()
overlay1.fill.fore_color.rgb = RGBColor(0x18, 0xCC, 0xB8)
overlay1.line.fill.background()

overlay2 = sl.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               cta_left, Inches(5.95), cta_w, Inches(0.2))
overlay2.fill.solid()
overlay2.fill.fore_color.rgb = RGBColor(0x10, 0xA0, 0x90)
overlay2.line.fill.background()

add_text_box(sl, cta_left + Inches(0.3), Inches(3.45), cta_w - Inches(0.6), Inches(0.3),
             "Contact", font_size=14, color=RGBColor(0xE0, 0xFF, 0xF0),
             alignment=PP_ALIGN.CENTER)
add_text_box(sl, cta_left + Inches(0.3), Inches(3.85), cta_w - Inches(0.6), Inches(0.5),
             "Bogdan Dr\u0103gu\u0219in", font_size=30, color=RGBColor(0xFF, 0xFF, 0xFF),
             bold=True, alignment=PP_ALIGN.CENTER)
add_text_box(sl, cta_left + Inches(0.3), Inches(4.5), cta_w - Inches(0.6), Inches(0.35),
             "contact@analize.online", font_size=20, color=RGBColor(0xFF, 0xFF, 0xFF),
             alignment=PP_ALIGN.CENTER)
add_text_box(sl, cta_left + Inches(0.3), Inches(5.15), cta_w - Inches(0.6), Inches(0.35),
             "Demo live:  analize.online", font_size=FONT_BODY, color=RGBColor(0xE0, 0xFF, 0xF0),
             alignment=PP_ALIGN.CENTER)

# Bottom accent line
add_gradient_bar(sl, Inches(3.5), Inches(6.0), Inches(6.3), TEAL_DIM, CYAN, TEAL, Pt(5))

add_text_box(sl, Inches(1), Inches(6.3), Inches(11.3), Inches(0.4),
             "Pre\u021bul nu va fi niciodat\u0103 mai mic dec\u00e2t ast\u0103zi.",
             font_size=17, color=AMBER, bold=True, alignment=PP_ALIGN.CENTER)

add_footer(sl, 14)

# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
# SAVE
# ━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━━
out_path = "C:/OldD/_Projects/Healthy/sale/AnalizeOnline-Pitch-RO.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
